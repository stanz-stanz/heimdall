"""Generate SIRI Startup Denmark pitch deck as .pptx.

Produces a 15-slide PowerPoint that can be imported into Canva for styling.
Content pulled from docs/business/heimdall-siri-application.md.

Run: python3 scripts/generate_pitch_deck.py
Output: data/output/heimdall-siri-pitch-deck.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def _add_slide(prs, title: str, bullets: list[str], notes: str = ""):
    """Add a slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.shapes.title.text = title

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.paragraphs[0].text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet

        para = tf.paragraphs[i]
        para.font.size = Pt(16)
        para.space_after = Pt(8)

        # Bold lines that start with specific markers
        if bullet.startswith("**") and bullet.endswith("**"):
            para.text = bullet.strip("*")
            para.font.bold = True
            para.font.size = Pt(18)
        elif "|" in bullet:
            para.font.size = Pt(14)
            para.font.name = "Courier New"

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def _add_title_slide(prs, title: str, subtitle: str):
    """Add the title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def _add_table_slide(prs, title: str, headers: list[str], rows: list[list[str]], notes: str = ""):
    """Add a slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    # Title
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    cols = len(headers)
    row_count = len(rows) + 1
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(0.4 * row_count)

    table = slide.shapes.add_table(row_count, cols, left, top, width, height).table

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True

    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- SLIDE 1: Cover ---
    _add_title_slide(
        prs,
        "Heimdall",
        "AI-Powered Cybersecurity for Small Businesses\n\n"
        "Federico Alvarez  |  Vejle, Denmark  |  2026\n\n"
        "Startup Denmark Application"
    )

    # --- SLIDE 2: The Problem ---
    _add_slide(prs, "The Problem", [
        "40% of Danish SMBs lack adequate cybersecurity (Styrelsen for Samfundssikkerhed)",
        "60% of SMBs that suffer a major breach close within 6 months (VikingCloud 2026)",
        "GDPR Article 32 requires security measures — most SMBs are non-compliant by default",
        "Every existing tool delivers through dashboards designed for security professionals",
        "The restaurant owner with outdated WordPress doesn't have a security team — she has Telegram",
    ], notes="The 40% gap is structural, not financial. The tools exist — they were never built for this audience.")

    # --- SLIDE 3: The Solution ---
    _add_slide(prs, "The Solution: Messaging-First Security", [
        "Continuous monitoring of public-facing digital surface (domains, certificates, CMS, plugins)",
        "Findings delivered as plain-language Telegram messages — not dashboards, not PDF reports",
        "AI-powered interpretation: Claude API translates technical findings into business language",
        "Persistent memory: Heimdall remembers what it told you and follows up on unresolved issues",
        "Shadow AI detection: first to scan for exposed AI agent infrastructure in SMBs",
    ], notes="No login portal. No dashboard. The owner reads it on the bus.")

    # --- SLIDE 4: Example Message ---
    _add_slide(prs, "What the Client Sees", [
        "**Heimdall Sikkerhedsadvarsel — uge 12**",
        "",
        "Vi fandt 2 ting pa jeres hjemmeside der kraever opmaerksomhed:",
        "",
        "1. Jeres hjemmeside korer pa en foraeldet version",
        "   Versionen I korer har 47 kendte sikkerhedshuller.",
        "",
        "2. Jeres SSL-certifikat udlober om 12 dage",
        "   Kunderne vil se en advarsel nar de prover at booke bord.",
    ], notes="Watchman tier example — plain language, no fix instructions, no jargon. In the client's preferred language.")

    # --- SLIDE 5: Innovation ---
    _add_slide(prs, "Six Innovations", [
        "1. Messaging-first delivery — the product IS the conversation, not a dashboard",
        "2. Digital Twin — CVE-level findings without touching client systems or needing consent",
        "3. Shadow AI detection — exposed AI agents are a new attack surface nobody else scans",
        "4. Persistent memory — longitudinal tracking creates switching costs",
        "5. Programmatic legal compliance (Valdi) — two-gate scanning validation with forensic logs",
        "6. AI interpretation chain — tools find, LLM explains, separation of detection from interpretation",
    ], notes="The digital twin is the key differentiator. §263 protects 'another person's data system' — the twin is ours.")

    # --- SLIDE 6: Market Opportunity ---
    _add_table_slide(prs, "Market Opportunity", [
        "Metric", "Calculation", "Annual Value"
    ], [
        ["TAM", "200,000 Danish SMBs with websites x 305 kr./mo x 12", "~732M kr./yr"],
        ["SAM", "80,000 (40% with inadequate security) x 305 kr./mo x 12", "~293M kr./yr"],
        ["SOM", "200 clients in 36 months x 305 kr./mo x 12", "~732K kr./yr"],
    ], notes="SOM is deliberately conservative: 200 clients = 0.25% of SAM. Upside: agency partnerships (1 relationship = 10-35 clients).")

    # --- SLIDE 7: Business Model ---
    _add_table_slide(prs, "Business Model & Pricing", [
        "Tier", "Monthly", "Annual", "Scanning", "Key Value"
    ], [
        ["Watchman (trial)", "199 kr.", "169 kr./mo", "Passive", "What is wrong, in plain language"],
        ["Sentinel", "399 kr.", "339 kr./mo", "Passive + Active", "What's wrong + how to fix it + daily monitoring"],
    ], notes="Two tiers: Watchman trial → Sentinel product. 399 kr./mo is cheaper than every competitor's comparable offering.")

    # --- SLIDE 8: Unit Economics ---
    _add_table_slide(prs, "Unit Economics (Per Client, Monthly)", [
        "Component", "Amount", "Notes"
    ], [
        ["Revenue (blended)", "~305 kr.", "Early mix, 30% annual uptake"],
        ["Claude API", "~50 kr.", "Interpretation + follow-up"],
        ["Infrastructure", "~15-30 kr.", "At 50+ clients"],
        ["Insurance", "~30-45 kr.", "Professional indemnity"],
        ["Total COGS", "~95-125 kr.", ""],
        ["Gross margin", "~59-69%", "Improves with scale"],
        ["Break-even", "~13-14 clients", "Fixed costs ~2,600 kr./mo"],
    ], notes="All prices excl. moms. Tool licensing eliminated (WPVulnerability API is free).")

    # --- SLIDE 9: Go-to-Market ---
    _add_slide(prs, "Go-to-Market Strategy", [
        "Phase 1 (Month 1-3): Vejle pilot — 5 clients via in-person visits, 'first finding free'",
        "Phase 2 (Month 3-6): Agency partnerships — '22 of your 35 client sites have issues'",
        "Phase 3 (Month 6-12): Local business networks — 20 clients, first operations hire",
        "Phase 4 (Month 12-24): Geographic expansion — Aarhus, Odense, Aalborg — 80 clients",
        "Phase 5 (Month 24-36): EU expansion — Germany, Netherlands — 200 clients, team of 6+",
        "",
        "Danish marketing law (Markedsforingslov) prohibits cold email — forces high-trust model",
        "This constraint becomes a moat: remote-first competitors can't replicate in-person trust",
    ], notes="One agency relationship = 10-35 client sites. The agency channel is highest leverage.")

    # --- SLIDE 10: Competitive Landscape ---
    _add_table_slide(prs, "Competitive Landscape", [
        "Competitor", "Starting Price", "Interface", "SMB Messaging"
    ], [
        ["Heimdall", "199 kr./mo", "Telegram/WhatsApp", "Yes"],
        ["Intruder.io", "~740 kr./mo", "Dashboard + Slack", "No"],
        ["Detectify", "~610 kr./mo", "Dashboard", "No"],
        ["HostedScan", "~215 kr./mo", "Dashboard + API", "No"],
        ["Beagle Security", "~885 kr./mo", "Dashboard", "No"],
    ], notes="Enterprise EASM players (CrowdStrike, Qualys) are moving upmarket, not down. They are not competing for the restaurant owner.")

    # --- SLIDE 11: Why Denmark ---
    _add_slide(prs, "Why Denmark", [
        "211M kr. government cybersecurity investment (2026-2029) — explicitly targets SMBs",
        "NCC-DK grants, Digital Europe Programme, Industriens Fond — accessible post-CVR",
        "GDPR-first: product built for strictest regime translates to all 27 EU states",
        "200,000+ SMBs with websites, high digital adoption — ideal first market",
        "Marketing law forces in-person model — being physically present is a competitive advantage",
        "I have been based in Vejle since 2019 — 7 years in Denmark, embedded in local business",
    ], notes="Denmark is not arbitrary. The government is creating the market; Heimdall serves it.")

    # --- SLIDE 12: Startup Phase Activities ---
    _add_slide(prs, "Startup Phase Plan", [
        "Month 1: Establish Heimdall ApS (CVR), open business bank account, professional indemnity insurance",
        "Month 1-2: Legal counsel — confirm §263 boundaries, scanning authorization template",
        "Month 1-3: Recruit 5 Vejle pilot clients (in-person, first finding free)",
        "Month 2-3: Pi5 production deployment, Telegram delivery pipeline live",
        "Month 3-6: Iterate on message tone, onboard agency partners",
        "Month 6: First operations hire (part-time, client communication)",
        "Month 6-12: Scale to 20 clients, local business network outreach",
    ], notes="Human-in-the-loop: I review every message before delivery during the pilot.")

    # --- SLIDE 13: Team & What's Built ---
    _add_slide(prs, "Team & Execution", [
        "Federico Alvarez — ~20 years enterprise software engineering (LEGO, JYSK, Deloitte, Medtronic)",
        "In Denmark since 2019, Vejle. Fast-Track employment scheme.",
        "Network security partner for domain expertise and technical credibility",
        "",
        "**What I have already built:**",
        "20+ module Python pipeline — tested against 204 Vejle domains in 8.5 minutes",
        "Full Telegram delivery bot — scan to client notification, operator approval flow",
        "Digital twin framework, Valdi compliance system, 690+ automated tests",
        "Built with Claude Code — solo developer operating at small team output",
    ], notes="This is not a slide deck. The product is running. The pipeline processes real domains.")

    # --- SLIDE 14: Financial Projections ---
    _add_table_slide(prs, "Financial Projections", [
        "", "Month 12", "Month 24", "Month 36"
    ], [
        ["Conservative — clients", "10", "50", "100"],
        ["Conservative — MRR", "3,050 kr.", "17,500 kr.", "37,000 kr."],
        ["Moderate — clients", "20", "80", "200"],
        ["Moderate — MRR", "6,100 kr.", "29,600 kr.", "74,000 kr."],
        ["Optimistic — clients", "30", "120", "300"],
        ["Optimistic — MRR", "9,150 kr.", "44,400 kr.", "111,000 kr."],
    ], notes="All scenarios use subscription revenue only. Tier migration (Watchman trial->Sentinel) is excluded upside. Self-sustaining on subscriptions alone.")

    # --- SLIDE 15: The Ask ---
    _add_slide(prs, "The Ask", [
        "**Startup Denmark residence permit to establish Heimdall ApS in Denmark**",
        "",
        "The product is built. The pipeline is tested. The market is ready.",
        "",
        "Denmark is investing 211M kr. to close the SMB cybersecurity gap.",
        "I am building the service that makes that investment reach the businesses that need it.",
        "",
        "199 kr./month. Telegram. Plain language. No dashboard required.",
    ], notes="Break-even at 13-14 clients. 68 prime targets identified in Vejle alone. I have proof of financial capacity for the establishment phase.")

    # Save
    out = Path("data/output/heimdall-siri-pitch-deck.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_deck()
